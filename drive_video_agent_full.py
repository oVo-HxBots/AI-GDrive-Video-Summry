#!/usr/bin/env python3
"""
drive_video_agent_full.py

Batch public-Google-Drive-video -> timestamped transcript -> structured report -> PPTX
Features:
 - Public Drive downloader (fragile but simple)
 - Optional Google Drive API download via service account (recommended for private files)
 - ffmpeg/ffprobe-based audio extraction and chunking (no pydub memory spikes)
 - Chunked transcription via OmniRoute/OpenAI-compatible /v1/audio/transcriptions
 - Timestamp-preserving segments when provider returns 'segments' (verbose_json)
 - Chat completions requesting structured JSON for final report + slides, with robust fallback
 - PPTX generation from structured JSON (or markdown fallback)
 - Manifest CSV + resume/retry logic for safe batch runs

Requirements:
 - ffmpeg and ffprobe in PATH
 - pip install requests python-pptx tqdm google-api-python-client google-auth
 - Set OMNIROUTE_API_BASE (for chat/completions) and OMNIROUTE_API_KEY in your environment. Optionally set OMNIROUTE_AUDIO_BASE to route /v1/audio requests to a different host (for example: http://161.118.182.88:8084) while using the same API key.
 - For Drive API downloads: set GOOGLE_SERVICE_ACCOUNT_JSON to the path of a service account JSON file on the machine running the script,
   and share the Drive files/folders with the service account email.

Usage:
    python drive_video_agent_full.py links.txt --workdir output --max-workers 3 --chunk-length 300
    python drive_video_agent_full.py links.txt --workdir output --resume --max-workers 2
"""
from __future__ import annotations
import os
import re
import sys
import time
import json
import csv
import logging
import argparse
import subprocess
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import List, Dict, Any, Optional
import requests
from pptx import Presentation
from pptx.util import Inches
from tqdm import tqdm

# Optional Google Drive API imports (only used when GOOGLE_SERVICE_ACCOUNT_JSON is set)
try:
    from google.oauth2 import service_account
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaIoBaseDownload
    GOOGLE_API_AVAILABLE = True
except Exception:
    GOOGLE_API_AVAILABLE = False

# ---------- config ----------
OMNI_BASE = os.getenv("OMNIROUTE_API_BASE", "").rstrip("/")
OMNI_AUDIO_BASE = os.getenv("OMNIROUTE_AUDIO_BASE", "").rstrip("/")  # optional separate base for /v1/audio endpoints
OMNI_KEY = os.getenv("OMNIROUTE_API_KEY")
MODEL_TRANSCRIBE = os.getenv("OMNIROUTE_MODEL_TRANSCRIBE", "whisper-1")
MODEL_CHAT = os.getenv("OMNIROUTE_MODEL_CHAT", "gpt-4o-mini")
MAX_RETRIES = 5
RETRY_BACKOFF_BASE = 1.5
DEFAULT_MAX_WORKERS = int(os.getenv("OMNI_MAX_WORKERS", "3"))
HEADERS = {"Authorization": f"Bearer {OMNI_KEY}"} if OMNI_KEY else {}

if not (OMNI_BASE or OMNI_AUDIO_BASE) or not OMNI_KEY:
    logging.warning("OMNIROUTE_API_BASE (chat) or OMNIROUTE_AUDIO_BASE (audio) or OMNIROUTE_API_KEY may not be set. HTTP calls will likely fail if required values are missing.")

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")


# ---------- utilities ----------
def get_drive_file_id(url_or_id: str) -> str:
    m = re.search(r"/d/([a-zA-Z0-9_-]+)", url_or_id)
    if m:
        return m.group(1)
    m = re.search(r"id=([a-zA-Z0-9_-]+)", url_or_id)
    if m:
        return m.group(1)
    # assume raw id
    if re.match(r"^[a-zA-Z0-9_-]{10,}$", url_or_id):
        return url_or_id
    raise ValueError(f"Could not parse Drive file ID from: {url_or_id}")


# ---------- Google Drive API (optional) ----------
def download_with_drive_api(file_id: str, dest_path: str) -> str:
    """
    Use a service account JSON specified by GOOGLE_SERVICE_ACCOUNT_JSON to download a file.
    Requires google-api-python-client and google-auth libraries. The service account or its domain
    must have access to the file (share the file/folder with the service account email).
    """
    sa_path = os.getenv("GOOGLE_SERVICE_ACCOUNT_JSON")
    if not sa_path:
        raise RuntimeError("GOOGLE_SERVICE_ACCOUNT_JSON not set; cannot use Drive API")
    if not GOOGLE_API_AVAILABLE:
        raise RuntimeError("google-api-python-client or google-auth not installed")
    scopes = ["https://www.googleapis.com/auth/drive.readonly"]
    creds = service_account.Credentials.from_service_account_file(sa_path, scopes=scopes)
    service = build("drive", "v3", credentials=creds, cache_discovery=False)
    meta = service.files().get(fileId=file_id, fields="id,name,mimeType,size").execute()
    request = service.files().get_media(fileId=file_id)
    with open(dest_path, "wb") as fh:
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
            if status:
                logging.debug(f"Download {int(status.progress() * 100)}%")
    return dest_path


# ---------- public Drive downloader (fragile) ----------
def download_public_drive_file(file_id: str, dest_path: str, chunk_size: int = 32768, timeout: int = 60) -> str:
    """
    Uses docs.google.com/uc?export=download to fetch public 'anyone with link' files.
    This method is fragile; consider using Drive API & OAuth for production.
    """
    url = "https://docs.google.com/uc?export=download"
    session = requests.Session()
    resp = session.get(url, params={"id": file_id}, stream=True, timeout=timeout)
    def get_confirm_token(r):
        for k, v in r.cookies.items():
            if k.startswith("download_warning"):
                return v
        return None
    token = get_confirm_token(resp)
    if token:
        resp = session.get(url, params={"id": file_id, "confirm": token}, stream=True, timeout=timeout)
    resp.raise_for_status()
    with open(dest_path, "wb") as f:
        for chunk in resp.iter_content(chunk_size):
            if chunk:
                f.write(chunk)
    return dest_path

def download_drive_file(file_id: str, dest_path: str) -> str:
    """
    Try Drive API if GOOGLE_SERVICE_ACCOUNT_JSON is set and google client libraries are available,
    otherwise fall back to public downloader.
    """
    if os.getenv("GOOGLE_SERVICE_ACCOUNT_JSON"):
        try:
            logging.info("Attempting download via Google Drive API (service account).")
            return download_with_drive_api(file_id, dest_path)
        except Exception as e:
            logging.warning(f"Drive API download failed ({e}); falling back to public downloader.")
    logging.info("Downloading via public Drive downloader (fragile).")
    return download_public_drive_file(file_id, dest_path)


# ---------- Run shell commands capturing stderr/stdout ----------
def run_cmd(cmd: List[str], raise_on_err: bool = True) -> (str, str):
    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    stdout = proc.stdout or ""
    stderr = proc.stderr or ""
    if proc.returncode != 0:
        logging.debug(f"CMD: {' '.join(cmd)}")
        logging.debug(f"STDOUT: {stdout}")
        logging.debug(f"STDERR: {stderr}")
        if raise_on_err:
            raise RuntimeError(f"Command failed: {' '.join(cmd)}\nStderr: {stderr.strip()}")
    return stdout, stderr


# ---------- ffmpeg-based audio extraction and splitting ----------
def extract_audio_ffmpeg(video_path: str, audio_out: str, sample_rate: int = 16000) -> str:
    cmd = [
        "ffmpeg", "-y", "-i", video_path,
        "-vn",
        "-acodec", "pcm_s16le",
        "-ar", str(sample_rate),
        "-ac", "1",
        audio_out
    ]
    run_cmd(cmd)
    return audio_out

def split_audio_ffmpeg(input_wav: str, out_dir: str, chunk_length_sec: int = 300) -> List[str]:
    Path(out_dir).mkdir(parents=True, exist_ok=True)
    out_pattern = os.path.join(out_dir, "chunk_%03d.wav")
    cmd = [
        "ffmpeg", "-y", "-i", input_wav,
        "-f", "segment",
        "-segment_time", str(chunk_length_sec),
        "-acodec", "pcm_s16le",
        "-ar", "16000",
        "-ac", "1",
        out_pattern
    ]
    run_cmd(cmd)
    files = sorted(Path(out_dir).glob("chunk_*.wav"))
    return [str(p) for p in files]


# ---------- media duration & formatting ----------
def get_media_duration(path: str) -> float:
    cmd = ["ffprobe", "-v", "error", "-show_entries", "format=duration", "-of", "json", path]
    out, err = run_cmd(cmd)
    try:
        data = json.loads(out)
        return float(data["format"]["duration"])
    except Exception:
        return 0.0

def format_time(seconds: float) -> str:
    seconds = int(round(seconds))
    h = seconds // 3600
    m = (seconds % 3600) // 60
    s = seconds % 60
    return f"{h:02d}:{m:02d}:{s:02d}"


# ---------- HTTP helpers with retries ----------
def http_post_with_retries(url: str, headers: Optional[dict] = None, files: Optional[dict] = None,
                           json_body: Optional[dict] = None, data: Optional[dict] = None, timeout: int = 300) -> requests.Response:
    headers = headers or {}
    last_exc = None
    for attempt in range(1, MAX_RETRIES + 1):
        try:
            if files:
                r = requests.post(url, headers=headers, files=files, data=data, timeout=timeout)
            else:
                r = requests.post(url, headers={**headers, "Content-Type": "application/json"}, json=json_body, timeout=timeout)
            if r.status_code in (429, 502, 503, 504):
                raise requests.HTTPError(f"Transient status {r.status_code}: {r.text}")
            r.raise_for_status()
            return r
        except Exception as e:
            last_exc = e
            backoff = (RETRY_BACKOFF_BASE ** attempt) + (0.1 * attempt)
            logging.warning(f"HTTP request failed (attempt {attempt}): {e}. Backing off {backoff:.1f}s")
            time.sleep(backoff)
    raise last_exc


# ---------- transcription (chunk-level) ----------
def transcribe_chunk_with_timestamps(chunk_path: str, chunk_start_sec: float, model: str = MODEL_TRANSCRIBE,
                                     request_verbose_json: bool = True) -> List[Dict[str, Any]]:
    """
    Returns list of segments: {"start": float, "end": float, "text": str}
    """
    # Use OMNIROUTE_AUDIO_BASE if provided, otherwise fall back to OMNIROUTE_API_BASE
    base_for_audio = OMNI_AUDIO_BASE or OMNI_BASE
    if not base_for_audio:
        raise RuntimeError("OMNIROUTE_API_BASE or OMNIROUTE_AUDIO_BASE not set")
    url = base_for_audio + "/v1/audio/transcriptions"
    with open(chunk_path, "rb") as f:
        files = {"file": (Path(chunk_path).name, f, "audio/wav")}
        data = {"model": model}
        if request_verbose_json:
            data["response_format"] = "verbose_json"
        r = http_post_with_retries(url, headers=HEADERS, files=files, data=data)
        resp = r.json()
    segments: List[Dict[str, Any]] = []
    if isinstance(resp, dict) and "segments" in resp and isinstance(resp["segments"], list):
        logging.debug("Transcription provider returned segments (with timestamps).")
        for seg in resp["segments"]:
            s_rel = float(seg.get("start", 0.0))
            e_rel = float(seg.get("end", s_rel))
            s = s_rel + chunk_start_sec
            e = e_rel + chunk_start_sec
            text = seg.get("text", "").strip()
            if text:
                segments.append({"start": s, "end": e, "text": text})
    else:
        logging.info("Transcription provider did not return 'segments'; using coarse chunk timestamp.")
        text = resp.get("text") or resp.get("transcription") or ""
        dur = get_media_duration(chunk_path) or 0.0
        segments.append({"start": float(chunk_start_sec), "end": float(chunk_start_sec + dur), "text": text.strip()})
    return segments


# ---------- chat completion helpers (structured JSON preference) ----------
def chat_complete_raw(messages: List[Dict[str, str]], model: str = MODEL_CHAT, temperature: float = 0.2, max_tokens: int = 1200) -> str:
    if not OMNI_BASE:
        raise RuntimeError("OMNIROUTE_API_BASE not set")
    url = OMNI_BASE + "/v1/chat/completions"
    payload = {"model": model, "messages": messages, "temperature": temperature, "max_tokens": max_tokens}
    r = http_post_with_retries(url, headers=HEADERS, json_body=payload)
    resp = r.json()
    try:
        return resp["choices"][0]["message"]["content"]
    except Exception:
        return resp.get("text", "")

def chat_complete_json_with_retries(messages: List[Dict[str, str]], model: str = MODEL_CHAT, temperature: float = 0.1, max_tokens: int = 1200, attempts: int = 3) -> Any:
    """
    Call the chat completion and attempt to parse JSON. If the model doesn't return valid JSON, re-prompt up to `attempts` times.
    Returns parsed JSON or {'_raw_text': raw} on final failure.
    """
    last_raw = None
    for attempt in range(1, attempts + 1):
        raw = chat_complete_raw(messages, model=model, temperature=temperature, max_tokens=max_tokens)
        last_raw = raw
        txt = raw.strip()
        try:
            return json.loads(txt)
        except Exception:
            import re
            m = re.search(r"(\{[\s\S]*\}|\[[\s\S]*\])", txt)
            if m:
                try:
                    return json.loads(m.group(1))
                except Exception:
                    logging.debug("Found JSON-like block but failed to parse on attempt %d", attempt)
        if attempt < attempts:
            # prepare a repair prompt
            repair_msg = (
                "The previous assistant reply was not valid JSON. Here is the reply:\n\n" + raw +
                "\n\nPlease output ONLY valid JSON that matches the requested schema and nothing else. If you cannot, reply with an empty JSON object {}."
            )
            messages = messages + [{"role": "user", "content": repair_msg}]
            logging.info("Re-prompting LLM for valid JSON (attempt %d/%d)", attempt+1, attempts)
            time.sleep(1 + attempt)
    return {"_raw_text": last_raw}

def chat_complete_json(messages: List[Dict[str, str]], model: str = MODEL_CHAT, temperature: float = 0.1, max_tokens: int = 1200) -> Any:
    # Backwards-compatible wrapper that uses the retrying JSON parser
    return chat_complete_json_with_retries(messages, model=model, temperature=temperature, max_tokens=max_tokens, attempts=3)


# ---------- summarization prompts (JSON-enforcing) ----------
def summarize_chunk_omni_json(chunk: str) -> Any:
    system_prompt = (
        "You are an assistant that summarizes timestamped transcripts. "
        "Output MUST be a JSON object with fields: 'summary' (string), 'entities' (list of strings), "
        "'suggested_slides' (list of objects {title, bullets:list[string], speaker_notes:string}), "
        "and 'timestamps' (list of strings)."
    )
    user_prompt = "Here is a timestamped transcript chunk. Produce the JSON object described in the system prompt.\n\n" + chunk
    messages = [{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}]
    return chat_complete_json(messages, model=MODEL_CHAT, temperature=0.1, max_tokens=800)

def combine_summaries_omni_json(summaries: List[Any]) -> Any:
    system = (
        "You are an assistant that produces internship-style project reports. "
        "Return a JSON object with keys: title, objectives, tools_methods, dataset_materials, key_results, conclusion, recommendations, slides (list). "
        "Each slide: title, bullets(list), speaker_notes(string)."
    )
    user = "Create the structured report JSON from these chunk summaries:\n\n" + json.dumps(summaries, ensure_ascii=False)
    messages = [{"role": "system", "content": system}, {"role": "user", "content": user}]
    return chat_complete_json(messages, model=MODEL_CHAT, temperature=0.1, max_tokens=1400)


# ---------- transcript builder ----------
def build_timestamped_transcript(segments: List[Dict[str, Any]]) -> (str, str):
    lines = []
    srt = []
    for idx, seg in enumerate(segments, start=1):
        start_s = format_time(seg["start"])
        end_s = format_time(seg["end"])
        text = seg["text"].replace("\n", " ").strip()
        lines.append(f"[{start_s}] {text}")
        srt.append(str(idx))
        srt.append(f"{start_s.replace('.',',')},000 --> {end_s.replace('.',',')},000")
        srt.append(text)
        srt.append("")
    return "\n".join(lines), "\n".join(srt)


# ---------- PPT generation ----------
def create_pptx_from_report_structured(report_obj: Dict[str, Any], thumbnails: List[str], out_pptx: str) -> str:
    prs = Presentation()
    title = report_obj.get("title") or (report_obj.get("slides", [{}])[0].get("title") if report_obj.get("slides") else "Project Summary")
    slide_layout = prs.slide_layouts[0]
    s0 = prs.slides.add_slide(slide_layout)
    s0.shapes.title.text = title[:200]
    slides = report_obj.get("slides") or report_obj.get("suggested_slides") or []
    for slide in slides[:20]:
        t = str(slide.get("title", "Slide"))[:90]
        bullets = slide.get("bullets", []) or slide.get("bullet_points", [])
        notes = slide.get("speaker_notes", "") or slide.get("notes", "")
        sl = prs.slides.add_slide(prs.slide_layouts[1])
        sl.shapes.title.text = t
        tf = sl.shapes.placeholders[1].text_frame
        for b in bullets:
            p = tf.add_paragraph(); p.text = str(b); p.level = 0
        if notes:
            try:
                ns = sl.notes_slide
                ns.notes_text_frame.text = notes[:4000]
            except Exception:
                pass
    # key frames
    if thumbnails:
        try:
            layout_idx = 5 if len(prs.slide_layouts) > 5 else 1
            s = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            try:
                s.shapes.title.text = "Key frames"
            except Exception:
                pass
            left = Inches(0.5); top = Inches(1.2)
            for thumb in thumbnails[:6]:
                try:
                    s.shapes.add_picture(thumb, left, top, width=Inches(2))
                except Exception:
                    pass
                left += Inches(2)
                if left > Inches(8):
                    left = Inches(0.5); top += Inches(1.5)
        except Exception:
            logging.warning("Failed to create key frames slide.")
    prs.save(out_pptx)
    return out_pptx

def create_pptx_from_report_fallback(markdown_text: str, thumbnails: List[str], out_pptx: str) -> str:
    prs = Presentation()
    first_line = markdown_text.splitlines()[0] if markdown_text else "Project Summary"
    s0 = prs.slides.add_slide(prs.slide_layouts[0]); s0.shapes.title.text = first_line[:200]
    sections = [p.strip() for p in markdown_text.split("\n\n") if p.strip()]
    for sec in sections[1:13]:
        lines = sec.splitlines()
        title = lines[0][:90] if lines else "Section"
        body = "\n".join(lines[1:])[:1500]
        sl = prs.slides.add_slide(prs.slide_layouts[1]); sl.shapes.title.text = title
        tf = sl.shapes.placeholders[1].text_frame
        for ln in body.splitlines():
            p = tf.add_paragraph(); p.text = ln.strip(); p.level = 0
    if thumbnails:
        try:
            s = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts)>5 else prs.slide_layouts[1])
            s.shapes.title.text = "Key frames"
            left = Inches(0.5); top = Inches(1.2)
            for thumb in thumbnails[:6]:
                try:
                    s.shapes.add_picture(thumb, left, top, width=Inches(2))
                except Exception:
                    pass
                left += Inches(2)
                if left > Inches(8):
                    left = Inches(0.5); top += Inches(1.5)
        except Exception:
            pass
    prs.save(out_pptx)
    return out_pptx


# ---------- thumbnails via ffmpeg ----------
def extract_thumbnails(video_path: str, out_dir: str, max_frames: int = 6) -> List[str]:
    Path(out_dir).mkdir(parents=True, exist_ok=True)
    duration = get_media_duration(video_path)
    if duration <= 0:
        return []
    times = [max(1, int(duration * i / (max_frames + 1))) for i in range(1, max_frames + 1)]
    out_files = []
    for idx, t in enumerate(times, start=1):
        outf = os.path.join(out_dir, f"thumb_{idx:03d}.jpg")
        cmd = ["ffmpeg", "-y", "-ss", str(t), "-i", video_path, "-frames:v", "1", "-q:v", "2", outf]
        try:
            run_cmd(cmd)
            if os.path.exists(outf):
                out_files.append(outf)
        except Exception as e:
            logging.warning(f"Thumbnail extraction failed for {video_path} at {t}s: {e}")
    return out_files


# ---------- per-file pipeline ----------
def process_drive_video(url_or_id: str, file_workdir: str, chunk_length_sec: int = 300,
                        request_verbose_json: bool = True, progress: Optional[tqdm] = None) -> Dict[str, Any]:
    res = {"input": url_or_id, "file_workdir": file_workdir, "success": False, "error": None}
    try:
        Path(file_workdir).mkdir(parents=True, exist_ok=True)
        file_id = get_drive_file_id(url_or_id)
        video_path = os.path.join(file_workdir, f"{file_id}.mp4")
        logging.info(f"Downloading {file_id} -> {video_path}")
        download_drive_file(file_id, video_path)
        # extract audio
        audio_path = os.path.join(file_workdir, f"{file_id}.wav")
        logging.info(f"Extracting audio to {audio_path}")
        extract_audio_ffmpeg(video_path, audio_path, sample_rate=16000)
        # split audio
        chunks_dir = os.path.join(file_workdir, "audio_chunks")
        logging.info(f"Splitting audio into chunks (len={chunk_length_sec}s) in {chunks_dir}")
        chunk_files = split_audio_ffmpeg(audio_path, chunks_dir, chunk_length_sec=chunk_length_sec)
        # transcribe chunks and accumulate segments
        segments = []
        current_start = 0.0
        for cp in chunk_files:
            dur = get_media_duration(cp)
            segs = transcribe_chunk_with_timestamps(cp, current_start, model=MODEL_TRANSCRIBE, request_verbose_json=request_verbose_json)
            segments.extend(segs)
            current_start += dur
            time.sleep(0.2)
        # build transcript and srt
        transcript_text, srt_text = build_timestamped_transcript(segments)
        transcript_path = os.path.join(file_workdir, f"{file_id}_transcript.txt")
        srt_path = os.path.join(file_workdir, f"{file_id}.srt")
        with open(transcript_path, "w", encoding="utf-8") as f: f.write(transcript_text)
        with open(srt_path, "w", encoding="utf-8") as f: f.write(srt_text)
        # summarize via LLM; request JSON per chunk and combine
        # chunk the transcript_text into manageable sizes
        def chunk_chars(text: str, n: int = 18000) -> List[str]:
            out=[]; i=0
            while i < len(text):
                out.append(text[i:i+n]); i+=n
            return out
        tr_chunks = chunk_chars(transcript_text, 18000)
        chunk_summaries = []
        for c in tr_chunks:
            summary = summarize_chunk_omni_json(c)
            chunk_summaries.append(summary)
            time.sleep(0.2)
        final_report_obj = combine_summaries_omni_json(chunk_summaries)
        # generate report + pptx
        thumbs_dir = os.path.join(file_workdir, "thumbs")
        thumbs = extract_thumbnails(video_path, thumbs_dir, max_frames=6)
        out_ppt = os.path.join(file_workdir, f"{file_id}_presentation.pptx")
        if isinstance(final_report_obj, dict) and "_raw_text" not in final_report_obj:
            # structured JSON report
            report_path = os.path.join(file_workdir, f"{file_id}_report.json")
            with open(report_path, "w", encoding="utf-8") as f:
                json.dump(final_report_obj, f, indent=2, ensure_ascii=False)
            create_pptx_from_report_structured(final_report_obj, thumbs, out_ppt)
        else:
            # fallback - LLM didn't provide structured JSON
            raw_md = final_report_obj.get("_raw_text") if isinstance(final_report_obj, dict) else str(final_report_obj)
            report_path = os.path.join(file_workdir, f"{file_id}_report.md")
            with open(report_path, "w", encoding="utf-8") as f:
                f.write(raw_md)
            create_pptx_from_report_fallback(raw_md, thumbs, out_ppt)
        res.update({
            "success": True,
            "video": video_path,
            "audio": audio_path,
            "transcript": transcript_path,
            "srt": srt_path,
            "report": report_path,
            "pptx": out_ppt
        })
    except Exception as e:
        logging.exception("Processing failed for %s", url_or_id)
        res.update({"error": str(e)})
    finally:
        if progress:
            progress.update(1)
        return res


# ---------- batch processing with manifest + resume/retry ----------
def read_links_file(path: str) -> List[str]:
    with open(path, "r", encoding="utf-8") as f:
        return [l.strip() for l in f.readlines() if l.strip() and not l.strip().startswith("#")]

def write_manifest_row(csv_path: str, row: Dict[str, Any], fieldnames: List[str]):
    exists = os.path.exists(csv_path)
    with open(csv_path, "a", newline="", encoding="utf-8") as fh:
        writer = csv.DictWriter(fh, fieldnames=fieldnames)
        if not exists:
            writer.writeheader()
        writer.writerow(row)

def load_manifest(csv_path: str) -> Dict[str, Dict[str, Any]]:
    if not os.path.exists(csv_path):
        return {}
    out = {}
    with open(csv_path, "r", encoding="utf-8") as fh:
        reader = csv.DictReader(fh)
        for r in reader:
            out[r.get("input") or r.get("file_id") or ""] = r
    return out

def process_multiple_links(links: List[str], workdir: str = "output", max_workers: int = DEFAULT_MAX_WORKERS,
                           manifest_name: str = "manifest.csv", resume: bool = False, retry_failed: bool = False,
                           force: bool = False, chunk_length_sec: int = 300) -> List[Dict[str, Any]]:
    Path(workdir).mkdir(parents=True, exist_ok=True)
    manifest_path = os.path.join(workdir, manifest_name)
    manifest_fieldnames = ["timestamp", "input", "file_id", "status", "video", "audio", "transcript", "srt", "report", "pptx", "error"]
    existing = load_manifest(manifest_path) if resume or retry_failed else {}
    # decide which links to process
    to_process = []
    for url in links:
        file_id = None
        try:
            file_id = get_drive_file_id(url)
        except Exception:
            file_id = url
        existing_entry = existing.get(url) or existing.get(file_id)
        if force:
            to_process.append((url, file_id))
            continue
        if resume and existing_entry and existing_entry.get("status") == "success":
            logging.info(f"Skipping {url} (already successful in manifest).")
            continue
        if retry_failed and existing_entry and existing_entry.get("status") == "failed":
            to_process.append((url, file_id))
            continue
        if existing_entry and existing_entry.get("status") == "success":
            logging.info(f"Skipping {url} (already successful in manifest).")
            continue
        # default: process
        to_process.append((url, file_id))
    results = []
    progress = tqdm(total=len(to_process), desc="Processing videos")
    with ThreadPoolExecutor(max_workers=max_workers) as ex:
        futures = {}
        for url, file_id in to_process:
            file_workdir = os.path.join(workdir, file_id or get_drive_file_id(url))
            fut = ex.submit(process_drive_video, url, file_workdir, chunk_length_sec, True, progress)
            futures[fut] = (url, file_id)
        for fut in as_completed(futures):
            url, file_id = futures[fut]
            res = fut.result()
            results.append(res)
            # write manifest row
            row = {
                "timestamp": time.strftime("%Y-%m-%d %H:%M:%S"),
                "input": url,
                "file_id": file_id or "",
                "status": "success" if res.get("success") else "failed",
                "video": res.get("video",""),
                "audio": res.get("audio",""),
                "transcript": res.get("transcript",""),
                "srt": res.get("srt",""),
                "report": res.get("report",""),
                "pptx": res.get("pptx",""),
                "error": res.get("error",""),
            }
            write_manifest_row(manifest_path, row, manifest_fieldnames)
    progress.close()
    return results


# ---------- CLI ----------
def main():
    parser = argparse.ArgumentParser(description="Batch Drive video -> transcript -> summary -> PPTX (OmniRoute/OpenAI-compatible)")
    parser.add_argument("links_file", help="Text file with Drive links or IDs, one per line")
    parser.add_argument("--workdir", default="output", help="Output workdir")
    parser.add_argument("--max-workers", type=int, default=DEFAULT_MAX_WORKERS, help="Concurrent workers")
    parser.add_argument("--chunk-length", type=int, default=300, help="Audio chunk length in seconds (default 300)")
    parser.add_argument("--manifest", default="manifest.csv", help="Manifest CSV filename inside workdir")
    parser.add_argument("--resume", action="store_true", help="Skip items already marked success in manifest")
    parser.add_argument("--retry-failed", action="store_true", help="Only retry items marked failed in manifest")
    parser.add_argument("--force", action="store_true", help="Force reprocess all items regardless of manifest")
    args = parser.parse_args()

    links = read_links_file(args.links_file)
    if not links:
        print("No links found in", args.links_file)
        sys.exit(1)
    print(f"Processing {len(links)} links -> {args.workdir} with {args.max_workers} workers. Manifest: {args.manifest}")
    res = process_multiple_links(links, workdir=args.workdir, max_workers=args.max_workers,
                                 manifest_name=args.manifest, resume=args.resume, retry_failed=args.retry_failed,
                                 force=args.force, chunk_length_sec=args.chunk_length)
    succeeded = sum(1 for r in res if r.get("success"))
    print(f"Done. {succeeded}/{len(res)} succeeded. Manifest at {os.path.join(args.workdir, args.manifest)}")

if __name__ == "__main__":
    main()
